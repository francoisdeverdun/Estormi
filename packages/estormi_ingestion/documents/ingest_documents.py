#!/usr/bin/env python3
"""
Walk iCloud Drive (or any directory) and ingest text-extractable documents
into the Estormi MCP server.

Usage:
    python3 ingest_documents.py [--root <dir>] [--mcp-url <url>] [--dry-run]

Supported formats: PDF, DOCX, ODT, PPTX, XLSX, TXT, MD, RTF, CSV, HTML, JSON, YAML
Skipped:           images, video, audio, binaries, files > MAX_FILE_MB
"""

import argparse
import asyncio
import csv
import os
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

from estormi_ingestion.shared.chunker import paragraph_chunks
from estormi_ingestion.shared.config import mcp_url
from estormi_ingestion.shared.emit import content_base_hash, post_chunks
from estormi_ingestion.shared.watermark import get_watermark, is_future_watermark, set_watermark
from memory_core.pii_filter import filter_pii, redact_code_secrets

# ── config ────────────────────────────────────────────────────────────────────
# DOCUMENTS_ROOT must be set explicitly (settings → `documents_root` →
# exported by `estormi_server/server/jobs.py`). Previously defaulted to the
# user's iCloud Drive root, which silently walked the user's entire
# cloud folder on first run. Now we require a deliberate pick.
ICLOUD_ROOT = (
    Path(os.environ["DOCUMENTS_ROOT"]).expanduser() if os.getenv("DOCUMENTS_ROOT") else None
)
DEFAULT_MCP = mcp_url()
MAX_FILE_MB = 20
CHUNK_SIZE = 1200
CHUNK_MIN = 80

SKIP_DIRS = {
    ".git",
    "__pycache__",
    ".DS_Store",
    "node_modules",
    ".Trash",
    "lost+found",
}
SKIP_EXTENSIONS = {
    # media
    ".jpg",
    ".jpeg",
    ".png",
    ".gif",
    ".webp",
    ".heic",
    ".heif",
    ".svg",
    ".ico",
    ".bmp",
    ".tiff",
    ".mp4",
    ".mov",
    ".avi",
    ".mkv",
    ".m4v",
    ".mp3",
    ".m4a",
    ".aac",
    ".wav",
    ".flac",
    # archives & binaries
    ".zip",
    ".tar",
    ".gz",
    ".dmg",
    ".pkg",
    ".app",
    ".exe",
    ".bin",
    # iCloud placeholders
    ".icloud",
    # CAD / specialised
    ".ifc",
    ".dwg",
    ".dxf",
    # certs
    ".p12",
    ".pem",
    ".cer",
    ".key",
}

# ── text extractors ───────────────────────────────────────────────────────────


def extract_pdf(path: Path) -> str:
    import pdfplumber

    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages.append(text)
    return "\n".join(pages)


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_odt(path: Path) -> str:
    from odf.opendocument import load
    from odf.text import P

    doc = load(str(path))
    paragraphs = []
    for elem in doc.text.getElementsByType(P):
        parts = []
        for node in elem.childNodes:
            if hasattr(node, "data"):
                parts.append(node.data)
        line = "".join(parts).strip()
        if line:
            paragraphs.append(line)
    return "\n".join(paragraphs)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(r.text for r in para.runs).strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(rows)


def extract_csv(path: Path) -> str:
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        return "\n".join("\t".join(row) for row in reader if any(row))


def extract_html(path: Path) -> str:
    import html as html_lib
    from html.parser import HTMLParser

    class _Extractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip = True

        def handle_endtag(self, tag):
            if tag in ("script", "style"):
                self._skip = False

        def handle_data(self, data):
            if not self._skip and data.strip():
                self.parts.append(data)

    raw = path.read_text(errors="ignore")
    parser = _Extractor()
    parser.feed(raw)
    return html_lib.unescape(" ".join(parser.parts))


def extract_rtf(path: Path) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(path.read_text(errors="ignore"))


def extract_plain(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".odt": extract_odt,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
    ".html": extract_html,
    ".htm": extract_html,
    ".rtf": extract_rtf,
    ".txt": extract_plain,
    ".md": extract_plain,
    ".markdown": extract_plain,
    ".json": extract_plain,
    ".yaml": extract_plain,
    ".yml": extract_plain,
}

# ── helpers ───────────────────────────────────────────────────────────────────


def ensure_downloaded(path: Path, timeout: int = 60) -> bool:
    """
    Trigger iCloud download if needed and wait up to `timeout` seconds.
    Returns True when the file is locally available.
    """
    import subprocess
    import time

    if path.suffix.lower() == ".icloud":
        return False  # placeholder file, not the real one

    # Kick iCloud download (no-op if already local)
    subprocess.run(["brctl", "download", str(path)], capture_output=True, timeout=5)

    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            size = path.stat().st_size
            # Check it's not still being downloaded (xattr present = in progress)
            result = subprocess.run(
                ["xattr", "-p", "com.apple.icloud.item.downloading", str(path)],
                capture_output=True,
            )
            if result.returncode != 0 and size > 0:
                return True  # no downloading xattr + has bytes = ready
        except OSError:
            pass
        time.sleep(2)

    return False


def clean(text: str) -> str:
    # Collapse intra-line whitespace runs but PRESERVE blank-line paragraph
    # breaks — paragraph_chunks splits on them, so flattening every newline
    # to a space would erase the structure that keeps unrelated sections in
    # separate chunks.
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


_FILENAME_DATE_RE = re.compile(r"^\s*(\d{4})[-_ ](\d{1,2})(?:[-_ ](\d{1,2}))?")


def extract_document_date(path: Path) -> str:
    """Best-effort document date.

    Priority:
      1. Leading YYYY-MM[-DD] in the filename (a common convention for
         scanned admin docs: '2024-03 Blood test results.pdf').
      2. File mtime (last write).

    Returns ISO date string or "" if neither is reliable.
    """
    m = _FILENAME_DATE_RE.match(path.stem)
    if m:
        y = int(m.group(1))
        mo = int(m.group(2))
        d = int(m.group(3)) if m.group(3) else 1
        if 1990 <= y <= 2100 and 1 <= mo <= 12 and 1 <= d <= 31:
            try:
                return datetime(y, mo, d, tzinfo=timezone.utc).date().isoformat()
            except ValueError:
                pass
    try:
        ts = path.stat().st_mtime
        return datetime.fromtimestamp(ts, tz=timezone.utc).date().isoformat()
    except OSError:
        return ""


def ingest_file(
    path: Path, mcp_url: str, headers: dict, dry_run: bool
) -> tuple[int, int, int, bool, bool]:
    """Return (chunks_ok, chunks_skipped, chunks_failed, transient, unreadable).

    Two distinct failure flags, because the walker treats them oppositely:

    * ``transient`` — the file could not be processed for a reason that may
      resolve on a later run (iCloud download timeout, an I/O error reading a
      file mid-sync). The caller must NOT advance the watermark past it,
      otherwise its mtime pre-dates the new watermark and it is skipped forever;
      leaving the watermark put retries the file next run.
    * ``unreadable`` — the file is deterministically unprocessable: an encrypted
      PDF, a corrupt office file, a parser that chokes on malformed bytes.
      Retrying never helps, so this must NOT block the watermark — otherwise one
      password-protected PDF pins the whole source forever, re-walking every
      file nightly and keeping the run red. It is logged + counted for
      visibility and re-attempted automatically only if the file's mtime changes
      (e.g. it later gets decrypted / re-saved). Distinguished from
      ``transient`` by error class: an :class:`OSError` is transient (I/O);
      anything else from the extractor is a parse/format failure, i.e.
      unreadable.

    A clean empty result (unsupported extension, oversized, too short) sets
    neither flag — those return (…, False, False) and the watermark advances.
    """
    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return 0, 0, 0, False, False

    if not ensure_downloaded(path, timeout=120):
        print(f"  [skip: iCloud download failed/timeout] {path.name}")
        return 0, 0, 0, True, False

    mb = path.stat().st_size / 1_048_576
    if mb > MAX_FILE_MB:
        print(f"  [skip: {mb:.1f}MB > {MAX_FILE_MB}MB] {path.name}")
        return 0, 0, 0, False, False

    try:
        text = clean(extractor(path))
    except OSError as e:
        # I/O failure (file vanished, read error mid-sync). May resolve next
        # run, so keep it transient and hold the watermark back.
        print(f"  [extract I/O error → retry next run] {path.name}: {e!r}")
        return 0, 0, 0, True, False
    except Exception as e:
        # Parse/format failure (encrypted PDF, corrupt office file, pdfminer
        # assertion). Deterministic: retrying never succeeds and would pin the
        # watermark forever. Skip as unreadable — counted + logged, non-blocking.
        print(f"  [unreadable → skipped] {path.name}: {e!r}")
        return 0, 0, 0, False, True

    if len(text) < 50:
        return 0, 0, 0, False, False

    # Documents pass through the same PII filter as messages/mail — bank
    # statements and scanned admin docs leak the same data shapes. Strip
    # machine secrets (AWS/GitHub/Stripe/Slack) first, then human PII —
    # pasted credentials land in PDFs and exported notes too.
    text = redact_code_secrets(text)
    text = filter_pii(text)

    chunks = paragraph_chunks(text, max_size=CHUNK_SIZE, min_size=CHUNK_MIN)

    # Document source_id = absolute path on disk; stable across renames of the
    # chunks themselves, and lets the server replace stale chunks when the
    # file is edited in place.
    source_id = str(path)
    base_hash = content_base_hash(source_id, text)
    doc_date = extract_document_date(path)

    def _log(idx: int, status: str) -> None:
        prefix = "[dry] " if dry_run else ""
        print(f"  {prefix}[{path.name}] chunk {idx}: {status}")

    ok, skipped, failed = post_chunks(
        "documents",
        source_id,
        chunks,
        mcp_url=mcp_url,
        title=path.name,
        url=str(path),
        date=doc_date,
        meta={"pii_filtered": True},
        base_hash=base_hash,
        headers=headers,
        dry_run=dry_run,
        on_result=_log,
    )
    return ok, skipped, failed, False, False


# ── main ──────────────────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--root",
        default=str(ICLOUD_ROOT) if ICLOUD_ROOT else None,
        help="Directory to walk (must be set explicitly via --root or DOCUMENTS_ROOT)",
    )
    parser.add_argument("--mcp-url", default=DEFAULT_MCP)
    parser.add_argument("--dry-run", action="store_true", help="Extract text but don't POST to MCP")
    args = parser.parse_args()

    if not args.root:
        print(
            "[documents] ERROR: no root path configured. Open the Documents "
            "manage modal in Estormi and pick a folder before running this "
            "ingest. (Setting key: `documents_root`; env var: `DOCUMENTS_ROOT`.)",
            file=sys.stderr,
        )
        sys.exit(2)

    # `.resolve()` follows symlinks so a tilde-rooted symlink can't escape to
    root = Path(args.root).expanduser().resolve()
    if not root.exists():
        print(f"ERROR: {root} does not exist")
        sys.exit(1)

    headers = {"Content-Type": "application/json"}

    force_full = os.getenv("FORCE_FULL", "").lower() in ("1", "true", "yes")
    last_ts, _ = asyncio.run(get_watermark("documents"))
    last_run = datetime.fromisoformat(last_ts) if last_ts and not force_full else None
    if last_run is not None and last_run.tzinfo is None:
        last_run = last_run.replace(tzinfo=timezone.utc)
    if force_full and last_ts:
        print("[docs] FORCE_FULL=1 — ignoring watermark, re-ingesting all files.")
    # Capture the walk-start time at the very start. Previously we rewound by
    # 60s as a safety margin, but the rewind is only useful when the start
    # was captured AFTER any pre-walk work (iCloud downloads, etc.); now that
    # we capture *before* every side-effect, the watermark is already
    # conservative and a file modified during the walk is naturally
    # re-examined next run.
    walk_started_at = datetime.now(timezone.utc)

    # Guard against a forward-skewed watermark written by a previous run whose
    # clock was ahead of real time (NTP drift, VM jump, manual change).  After
    # the clock corrects, last_run would be in the future relative to every
    # file's mtime, permanently skipping them until real time catches up —
    # potentially years of data loss.  Detecting and resetting last_run to None
    # forces a full rescan; content_hash dedup on the server makes re-ingest
    # cheap and idempotent, and any file modified during the skew gap is picked
    # up instead of being lost.
    if is_future_watermark(last_run, walk_started_at):
        print(
            f"[docs] WARNING: stored watermark ({last_run.isoformat()}) is in the "
            "future — likely from a run under forward clock-skew. Ignoring it and "
            "rescanning all files. Content-hash dedup prevents duplicate ingestion."
        )
        last_run = None

    # Trigger iCloud download for all placeholder files before walking.
    # Placeholders are named .RealFile.pdf.icloud and are created when macOS
    # evicts files to save disk space (common at 3 AM with display off).
    # Fire-and-forget: `brctl download` only *requests* the download from the
    # iCloud daemon, so there is nothing to wait for — blocking on each call
    # serialized what the daemon happily parallelizes (an evicted library of
    # hundreds of files used to stall the walk for minutes before it began).
    # The per-file `ensure_downloaded` during the walk does the actual waiting.
    import subprocess as _sp

    # followlinks=False to match the ingest walk below: a symlink loop must not
    # hang this pre-pass, and we never request iCloud downloads outside `root`.
    for dp, _, fnames in os.walk(root, followlinks=False):
        for fn in fnames:
            if fn.startswith(".") and fn.endswith(".icloud"):
                real = Path(dp) / fn[1 : -len(".icloud")]
                _sp.Popen(
                    ["brctl", "download", str(real)],
                    stdout=_sp.DEVNULL,
                    stderr=_sp.DEVNULL,
                )

    print(f"[docs] Walking {root} …")
    total_ok = total_skipped = total_failed = total_files = total_watermark_skipped = 0
    total_transient = 0
    total_unreadable = 0

    # followlinks=False: a symlink loop would hang the walk, and a symlink
    # pointing outside `root` would exfiltrate arbitrary files into the index.
    for dirpath, dirnames, filenames in os.walk(root, followlinks=False):
        dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS and not d.startswith(".")]
        for fname in filenames:
            # iCloud evicts files as hidden placeholders: .RealName.ext.icloud
            # Remap to the real path so ensure_downloaded can fetch it.
            if fname.startswith(".") and fname.endswith(".icloud"):
                real_name = fname[1 : -len(".icloud")]
                path = Path(dirpath) / real_name
            elif fname.startswith("."):
                continue
            else:
                path = Path(dirpath) / fname
            ext = path.suffix.lower()
            if ext in SKIP_EXTENSIONS:
                continue
            if last_run is not None:
                try:
                    # iCloud sync sometimes writes a downloaded file with the
                    # server's mtime, which can predate our watermark — using
                    # mtime alone would permanently skip files whose content
                    # was edited on iOS. ``ctime`` (inode change time) is
                    # updated on every local write on macOS, so taking the
                    # max of the two correctly detects "new on this device".
                    st = path.stat()
                    effective_ts = max(st.st_mtime, st.st_ctime)
                    if datetime.fromtimestamp(effective_ts, tz=timezone.utc) <= last_run:
                        total_watermark_skipped += 1
                        continue
                except OSError:
                    pass
            total_files += 1
            ok, skipped, failed, transient, unreadable = ingest_file(
                path, args.mcp_url, headers, args.dry_run
            )
            total_ok += ok
            total_skipped += skipped
            total_failed += failed
            total_transient += int(transient)
            total_unreadable += int(unreadable)

    if last_run and total_watermark_skipped > 0 and total_files == 0:
        print(
            f"\n[docs] NOTE: all {total_watermark_skipped} files pre-date watermark "
            f"({last_run.date()}). Run FORCE_FULL=1 to re-ingest everything."
        )
    print(
        f"\n[docs] Done — {total_files} files scanned "
        f"({total_watermark_skipped} skipped by watermark), "
        f"{total_ok} chunks ingested, {total_skipped} duplicates skipped, "
        f"{total_failed} chunk POSTs failed, "
        f"{total_unreadable} unreadable (encrypted/corrupt, skipped)."
    )
    if args.dry_run:
        return
    if total_failed > 0 or total_transient > 0:
        # Advancing the watermark on a run with failed POSTs or *transient*
        # per-file failures (iCloud download timeout, I/O read error) silently
        # drops the affected files forever: their mtime pre-dates the new
        # watermark, so the next run's effective-ts check skips them. A
        # transient failure returns (0,0,0) — indistinguishable from a clean
        # empty result by the chunk counters alone — so it is tracked
        # separately. Leave the watermark untouched so the same files come back
        # next run for another attempt. (Deterministically *unreadable* files —
        # encrypted/corrupt — are NOT counted here: retrying never helps, so
        # they must not pin the watermark. See ingest_file.)
        print(
            f"[docs] Watermark NOT advanced — {total_failed} chunk(s) failed to POST, "
            f"{total_transient} file(s) failed transiently (download/extract). "
            "Affected files will be retried on the next run."
        )
        sys.exit(1)
    # walk_started_at was captured BEFORE the walk, so it already
    # pre-dates any file mtime updated during ingestion. No rewind
    # needed — concurrent edits are picked up on the next run.
    asyncio.run(set_watermark("documents", walk_started_at.isoformat()))


if __name__ == "__main__":
    main()
